import os
import hashlib
from typing import Tuple, Optional
import pypdf
import pdfplumber
from docx import Document
from pptx import Presentation
import pandas as pd
from io import BytesIO

def calculate_file_hash(file_path: str) -> str:
    """Calculate SHA-256 hash of a file."""
    sha256_hash = hashlib.sha256()
    with open(file_path, "rb") as f:
        for byte_block in iter(lambda: f.read(4096), b""):
            sha256_hash.update(byte_block)
    return sha256_hash.hexdigest()

def extract_text_from_pdf(file_path: str) -> Tuple[str, dict]:
    """
    Extract text from PDF file.
    Returns: (text, metadata) where metadata contains page numbers.
    """
    text_parts = []
    metadata = {"pages": []}
    
    try:
        # Try pdfplumber first (better for complex PDFs)
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text() or ""
                text_parts.append(f"--- Page {i} ---\n{page_text}\n")
                metadata["pages"].append({
                    "page_number": i,
                    "text_length": len(page_text)
                })
    except Exception as e:
        # Fallback to pypdf
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                for i, page in enumerate(pdf_reader.pages, 1):
                    page_text = page.extract_text() or ""
                    text_parts.append(f"--- Page {i} ---\n{page_text}\n")
                    metadata["pages"].append({
                        "page_number": i,
                        "text_length": len(page_text)
                    })
        except Exception as e2:
            raise Exception(f"Failed to extract PDF text: {str(e2)}")
    
    return "\n".join(text_parts), metadata

def extract_text_from_docx(file_path: str) -> Tuple[str, dict]:
    """Extract text from DOCX file."""
    doc = Document(file_path)
    text_parts = []
    
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)
    
    # Also extract from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join([cell.text for cell in row.cells])
            if row_text.strip():
                text_parts.append(row_text)
    
    metadata = {"paragraphs": len(doc.paragraphs), "tables": len(doc.tables)}
    return "\n".join(text_parts), metadata

def extract_text_from_pptx(file_path: str) -> Tuple[str, dict]:
    """Extract text from PPTX file."""
    prs = Presentation(file_path)
    text_parts = []
    slide_count = 0
    
    for slide in prs.slides:
        slide_count += 1
        slide_text = []
        
        # Extract from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        
        # Extract from notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                slide_text.append(f"[Notes]: {notes_slide.notes_text_frame.text}")
        
        if slide_text:
            text_parts.append(f"--- Slide {slide_count} ---\n" + "\n".join(slide_text) + "\n")
    
    metadata = {"slides": slide_count}
    return "\n".join(text_parts), metadata

def extract_text_from_excel(file_path: str) -> Tuple[str, dict]:
    """Extract text from Excel file (XLSX or XLS)."""
    text_parts = []
    sheet_names = []
    
    try:
        # Try XLSX first
        excel_file = pd.ExcelFile(file_path, engine='openpyxl')
    except:
        try:
            # Fallback to XLS
            excel_file = pd.ExcelFile(file_path, engine='xlrd')
        except Exception as e:
            raise Exception(f"Failed to read Excel file: {str(e)}")
    
    for sheet_name in excel_file.sheet_names:
        sheet_names.append(sheet_name)
        df = pd.read_excel(excel_file, sheet_name=sheet_name)
        
        # Convert to markdown table format
        text_parts.append(f"--- Sheet: {sheet_name} ---\n")
        text_parts.append(df.to_markdown(index=False))
        text_parts.append("\n")
    
    metadata = {"sheets": sheet_names, "sheet_count": len(sheet_names)}
    return "\n".join(text_parts), metadata

def extract_text_from_csv(file_path: str) -> Tuple[str, dict]:
    """Extract text from CSV file."""
    df = pd.read_csv(file_path)
    metadata = {"rows": len(df), "columns": list(df.columns)}
    return df.to_markdown(index=False), metadata

def extract_text_from_file(file_path: str, file_extension: str) -> Tuple[str, dict]:
    """
    Extract text from file based on extension.
    Returns: (text, metadata)
    """
    ext = file_extension.lower()
    
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext == '.docx':
        return extract_text_from_docx(file_path)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_path)
    elif ext in ['.xlsx', '.xls']:
        return extract_text_from_excel(file_path)
    elif ext == '.csv':
        return extract_text_from_csv(file_path)
    elif ext in ['.txt', '.md']:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        return text, {"type": "text"}
    else:
        raise ValueError(f"Unsupported file type: {ext}")

